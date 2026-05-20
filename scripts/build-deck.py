"""Generate the Claude Architect Foundations PPTX from the reference template.

Inherits layouts/theme from context-engineering-april-2026.pptx (Tim Warner's
proven O'Reilly live-training deck), then populates the deck against the
RICH layout library in Master 2 - Section Header, Demo, Course Title,
Important Statement, Comparison, References, Up Next, Module Overview, etc.

This is the 2026-05-21 rebuild. The previous version used only 4 layouts
(Title Slide, Title and Content, Title Only, Blank) and rendered demo/quote
slides as raw text boxes against unstyled Blank backgrounds. The rebuild
maps each slide role to a purpose-built Master 2 layout.

Slide budget (matches the current 4-segment + Segment 2.5 course flow):

- 4 intro slides:        cover, bio, course flow, opening epigraph
- 4 x ~10 segment slides: divider + content + demo + exercise + takeaways
- 2 Segment 2.5 callouts: bridge + self-study deep-dive pointer
- 5 outro slides:        recap, CCA-F brief, study path, resources, thanks

Output: C:/github/claude-architect/slides/warner-claude-architect-may-2026.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REF = Path(r"C:/github/context-engineering/instructor/context-engineering-april-2026.pptx")
OUT = Path(r"C:/github/claude-architect/slides/warner-claude-architect-may-2026.pptx")

# ---------------------------------------------------------------------------
# Layout vocabulary
#
# The reference template ships 81 layouts across 3 masters. Master 2 is the
# rich Pluralsight-style library; we use it for almost everything. See
# `_spikes/inspect_decks.py` and `_spikes/deck_inventory.txt` for the full
# inventory. We address layouts by exact name since some names appear in
# multiple masters (e.g. "Section Header" in M0/L2 vs M2/L1 vs "Section
# Header" - the latter is the one we want with a stronger visual identity).
# ---------------------------------------------------------------------------

# Master 0 (basic Office layouts) - used only for the cover/title slide
# because Master 2's "Course Title" layout has no TITLE placeholder and
# python-pptx's slide.shapes.title fallback fails on it.
LAYOUT_COVER = "Title Slide"                  # M0/L0  cover + subtitle

# Master 2 (rich library) - the workhorses
LAYOUT_CONTENT = "04_Content Slide (top title bar)"   # M2/L60  primary bullet workhorse
LAYOUT_SECTION = "Section Header"             # appears in M0/L2 and M2/L1; we'll match the rich one
LAYOUT_DEMO = "Demo"                          # M2/L4   one big BODY, 6 decorative shapes
LAYOUT_UPNEXT = "Up Next"                     # M2/L2   one BODY, decorative banner
LAYOUT_STATEMENT = "Important Statement"      # M2/L49  big-centered-title quote
LAYOUT_QUOTE_SHORT = "Quotation-short"        # M2/L50  TITLE + BODY for attribution
LAYOUT_COMPARISON = "Comparison: Point-by-Point"  # M2/L47  4 BODY placeholders for A/B tables
LAYOUT_SUMMARY = "Module Overview/Summary"    # M2/L3   2 BODY placeholders (per-section takeaways)
LAYOUT_REFS = "References-4 Items"            # M2/L54  TITLE + 4 BODY + pictures (resources)
LAYOUT_DEFINITION = "Definition"              # M2/L52  TITLE + 2 BODY
LAYOUT_THREE_ITEMS = "Three Item List with Icons"  # M2/L19 for 3-tier breakdowns

# Fallbacks
LAYOUT_TITLE_ONLY = "Title Only"
LAYOUT_BLANK = "Blank"

# ---------------------------------------------------------------------------
# Slide content
#
# Each entry: (layout_name, title, body_lines_or_None, speaker_notes_or_None)
# The body_lines list is treated as one paragraph per element. Empty strings
# become visual spacer paragraphs.
#
# Speaker notes follow Tim's voice rubric (FRAMER + bracketed echoes + exam
# tip) without inline em dashes, AWS, "ask" as noun, or glazing openers.
# ---------------------------------------------------------------------------

SLIDES: list[tuple[str, str, list[str] | None, str | None]] = [
    # ======================================================================
    # COURSE INTRO
    # ======================================================================
    (
        LAYOUT_COVER,
        "Claude Architect Foundations",
        ["A 4-hour O'Reilly Live Training", "Tim Warner  |  techtrainertim.com"],
        (
            "Cold open. This is a skills-first production-architecture course. "
            "Mention CCA-F (Anthropic's certification) but do NOT teach to exam items. "
            "Set expectations: heavy on live demos, light on slides. The notebooks ARE the course; "
            "slides are the scaffolding."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Tim Warner",
        [
            "Microsoft MVP (Azure AI)",
            "Pluralsight Principal Author - 200+ courses, 1M+ learners",
            "O'Reilly Live Training instructor",
            "Microsoft Press / Pearson senior content developer",
            "28+ years on the Microsoft stack",
            "",
            "techtrainertim.com  |  github.com/timothywarner-org",
        ],
        (
            "30-second bio. I teach from first principles. No glazing, no fluff. "
            "If they want production patterns and live demos, they're in the right room. "
            "Skip the resume part if pacing is tight; the website is the call-to-action."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Course flow - 4 segments x 50 min",
        [
            "Segment 1 - Building AI Agents That Use Tools",
            "Segment 2 - Tool Design, Integration, and Claude Code Workflows",
            "Segment 3 - Structured Output, Context, and Production Reliability",
            "Segment 4 - CCA-F Certification Capstone",
            "",
            "Plus self-study: Segment 2.5 - Control Surfaces and the Console Asset Surface",
            "",
            "10-minute breaks between segments. Total 4 hours.",
        ],
        (
            "Tell them where the slides and code live: github.com/timothywarner-org/claude-architect. "
            "Call out Segment 2.5 as the self-study deep dive - not on the 4-hour clock, but the "
            "highest exam-objective density in the repo. Cert briefing is the capstone (Segment 4), "
            "not the lead. We earn the cert talk after the skills."
        ),
    ),
    (
        LAYOUT_STATEMENT,
        '"The model is non-deterministic. Your architecture cannot be."',
        None,
        (
            "Frame the whole course in one sentence. Hooks are deterministic. Schemas are deterministic. "
            "Application-layer intercepts are deterministic. The model isn't, and that's fine - it's the "
            "right tool for ambiguous decisions. Production patterns build a deterministic shell around "
            "a non-deterministic kernel. Every demo in this course is an instance of that pattern."
        ),
    ),

    # ======================================================================
    # SEGMENT 1: AGENTIC ARCHITECTURE (Domain 1, 27%)
    # ======================================================================
    (
        LAYOUT_SECTION,
        "Segment 1: Building AI Agents That Use Tools",
        None,
        (
            "Domain 1 - Agentic Architecture and Orchestration. 50 minutes. "
            "27% of the CCA-F blueprint, the single heaviest domain. "
            "Live build: customer support agent + PreToolUse hook + coordinator-subagent demo."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "What you'll build",
        [
            "A customer support agent with 4 tools (lookup + refund + escalate)",
            "Hook-enforced policy - the $500 refund cap can't be talked around",
            "Structured escalation to human review with a typed summary",
            "Coordinator + subagent isolation demo",
            "",
            "By the end of the segment: you can sketch any agent topology",
        ],
        (
            "This is the live build. Set expectations: we WILL hit edge cases - that's the point. "
            "The model is going to TRY to refund more than the cap. The hook is going to say no. "
            "The aha moment is watching the model re-plan and escalate correctly."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "The agentic loop",
        [
            "1. Send request - model returns content + stop_reason",
            "2. Read stop_reason - branch on the value, never on prose",
            "3. tool_use? Execute tool, append tool_result, loop",
            "4. end_turn? Return to user, done",
            "5. pause_turn? Resume in the next request, no state change",
            "",
            "Six stop_reason values: end_turn, max_tokens, stop_sequence, tool_use, pause_turn, refusal",
        ],
        (
            "The hard rule: never parse natural language to detect completion. Always branch on the "
            "API field. Demo segment-1 prints [iter N] stop_reason= on every iteration so learners "
            "see the state transitions live. Refusal is rare; pause_turn shows up with long-running "
            "server tools like web_search."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Coordinator + subagents",
        [
            "ONE coordinator holds the user-facing thread",
            "Subagents get ISOLATED contexts via the Task tool",
            "Subagent returns its FINAL answer to the parent, not its working notes",
            "Split when: own success criteria, own tool subset, own context discipline",
            "",
            "Don't split for splitting's sake - two coordinators chatting is one slow agent",
        ],
        (
            "Isolation is the feature. The parent doesn't drown in subagent reasoning, the parent "
            "sees a clean return value. Reference cookbook: "
            "claude-cookbooks-main/claude_agent_sdk/01_The_chief_of_staff_agent.ipynb. "
            "Segment 1 demo notebook contains a stripped-down coordinator-subagent at the end."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Hooks - deterministic guarantees",
        [
            "PreToolUse - gate a call BEFORE it runs (the bouncer)",
            "PostToolUse - audit, transform, log AFTER (the auditor)",
            "SessionStart - inject context (the briefing room)",
            "Stop - final verification before returning to user",
            "",
            "Hooks run as YOUR code, not the model's. Anything that MUST happen, hook it.",
        ],
        (
            "The prompt advises; the hook enforces. The $500 refund cap is the canonical example: "
            "prompt says 'don't exceed cap', tool description says 'cap is 500', and the hook "
            "PHYSICALLY blocks anything over. When all three layers agree, the hook never fires - "
            "but it's there for when they disagree."
        ),
    ),
    (
        LAYOUT_DEMO,
        "Demo - customer support agent (15 min)",
        [
            "Open notebooks/segment-1-customer-support-agent.ipynb",
            "Wire 4 tools, run Scenario A ($80 refund, in-cap, agent cooperates)",
            "Add the PreToolUse hook from hooks-example.py",
            "Run Scenario B ($750 refund, over-cap, hook blocks, agent escalates)",
            "Coordinator-subagent demo at the bottom of the notebook",
        ],
        (
            "Live demo, 15 minutes. The aha moment is when the hook fires and the model "
            "re-plans to call escalate_to_human with ticket ESC-9001. Have a backup screenshot "
            "ready in case the API hiccups. Cookbook reference for further study: "
            "claude-cookbooks-main/tool_use/customer_service_agent.ipynb."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Exercise (5 min)",
        [
            "Sketch architecture for a Developer Productivity scenario",
            "On paper or in chat: coordinator + 3 subagents + tool scope per agent",
            "What does each subagent specialize in?",
            "What is each subagent's success criterion?",
            "What tools does each one get? Cap each at ~5.",
        ],
        (
            "Walk the chat for 90 seconds, then share 2-3 designs out loud. Reinforce: "
            "tool descriptions beat tool names, isolation beats shared context, ~5 tools per "
            "agent is the floor before selection accuracy degrades."
        ),
    ),
    (
        LAYOUT_SUMMARY,
        "Segment 1 takeaways",
        [
            "Branch on stop_reason, never on natural-language signals",
            "Coordinator + isolated subagents beats one giant context",
            "Hooks turn 'should' into 'must' - the prompt advises, the hook enforces",
        ],
        (
            "Bridge: 'You just built an agent that decides what to do. Next we go one level deeper - "
            "into the tools themselves and the Claude Code surface where you author them on a real team.'"
        ),
    ),

    # ======================================================================
    # SEGMENT 2: TOOLS + CLAUDE CODE (Domains 2 + 3, 38%)
    # ======================================================================
    (
        LAYOUT_SECTION,
        "Segment 2: Tool Design, Integration, Claude Code",
        None,
        (
            "Domains 2 + 3. 50 minutes. Combined 38% of exam weight - the biggest pairing on the blueprint. "
            "Two demos: MCP config walkthrough, claude -p headless mode."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Two domains, one segment",
        [
            "Domain 2 - Tool Design and MCP Integration (18%)",
            "Domain 3 - Claude Code Configuration and Workflows (20%)",
            "Combined = 38% of the cert, the biggest single chunk",
            "",
            "Why combined? Claude Code IS the primary MCP-consuming surface",
            "You cannot teach one without the other in production",
        ],
        (
            "Explain the merge. Tools are the verbs; Claude Code is where you orchestrate them. "
            "Different domains in the blueprint, same conversation in real life. Two demos "
            "in this segment - MCP config walkthrough first, then claude -p headless second."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Tool definitions - name + description + input_schema",
        [
            "name        - how the model invokes the tool",
            "description - STRONGLY recommended, drives tool selection",
            "input_schema - JSON Schema draft 2020-12 (type, properties, required)",
            "Optional: cache_control ephemeral (5m / 1h), strict: true",
            "",
            "Tool DESCRIPTIONS do the work, not tool NAMES",
        ],
        (
            "A tool named get_data with description 'Fetch customer record by ID, returns name, "
            "email, status, lifetime value' beats one named getCustomerRecordById with description "
            "'Gets data.' Write descriptions like you're onboarding a junior engineer. The model "
            "reads descriptions, not your variable-naming convention."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Structured tool errors",
        [
            "Anti-pattern: raise exceptions from tool implementations",
            "Anti-pattern: return empty strings on failure",
            "Pattern: return tool_result with is_error: true",
            "Add structured fields: errorCategory, isRetryable",
            "",
            "The model can then DECIDE: retry, fall back, or escalate",
        ],
        (
            "Thrown exceptions crash the agent. Structured errors let it recover. This is the entire "
            "reliability story for tools in one slide. The three categories we use: transient (retry "
            "with backoff), permanent (reformulate or surface), policy (escalate or change approach)."
        ),
    ),
    (
        LAYOUT_COMPARISON,
        "tool_choice modes",
        [
            "auto",
            "Model decides whether to call a tool. Default. May answer in prose. "
            "Pick when the model's judgment is the right call.",
            "any",
            "Model MUST call some tool, model picks which. Use when action is required.",
            "tool",
            "Model MUST call THIS specific tool. The structured-output cheat code. "
            "{type: tool, name: extract_X}.",
            "none",
            "Tools registered but model CANNOT call them. 'Explain, don't act' turns.",
        ],
        (
            "The traffic-light analogy from Segment 2.5: auto is green, any is green with a tow truck "
            "behind you, tool=X is a forced left turn, none is a red light. "
            "disable_parallel_tool_use lives INSIDE the tool_choice object, not as a top-level kwarg. "
            "Common multiple-choice trap."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Cap agents at ~5 tools",
        [
            "More tools = worse selection accuracy (well-documented effect)",
            "Beyond 5: split into subagents with scoped tool sets",
            "Specificity in descriptions beats quantity of tools",
            "",
            "Anti-pattern: dump 30 tools into one agent because 'flexibility'",
        ],
        (
            "Mention the analyze_dependencies vs Grep anti-pattern: model defaults to built-in Grep "
            "if the custom tool description isn't specific enough. Description is the contract. "
            "Tool count is the soft cap."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "MCP server config - .mcp.json",
        [
            "Three transports: stdio, SSE, HTTP",
            "stdio: command + args + env (local subprocess)",
            "SSE / HTTP: type + url + headers (remote)",
            "${ENV_VAR} expansion works in env, args, AND headers",
            "",
            "Project scope: ./.mcp.json    User scope: ~/.claude.json",
        ],
        (
            "Show repo's own .mcp.json in the demo. Four servers, three transports, env-var expansion "
            "for secrets. The pattern that scales: add a new MCP server to .mcp.json, every MCP-aware "
            "agent picks it up on the next list_tools call - no source change."
        ),
    ),
    (
        LAYOUT_DEMO,
        "Demo A - MCP config walkthrough (10 min)",
        [
            "Open C:/github/claude-architect/.mcp.json in VS Code",
            "Walk server-by-server: filesystem (stdio), github (stdio + ${GITHUB_TOKEN})",
            "Show context7 (HTTP + header), internal-knowledge-base (SSE + bearer)",
            "Show what happens when GITHUB_TOKEN is unset",
        ],
        (
            "Live demo, 10 minutes. Open .mcp.json side-by-side with the notebook so learners see "
            "config-as-data. Mention that the file is the same shape Claude Code reads at startup - "
            "it's the SDK contract, not a Claude Code convention."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "CLAUDE.md hierarchy",
        [
            "User     ~/.claude/CLAUDE.md     - your machine, every project",
            "Project  ./CLAUDE.md             - team conventions, in git",
            "Subtree  ./<subdir>/CLAUDE.md    - loads on demand when files in subdir are read",
            "Local    ./CLAUDE.local.md       - gitignored personal override",
            "",
            "Agent SDK: settingSources [user, project] loads both",
        ],
        (
            "Load order matters. User and project are additive. Subtree on-demand keeps the token "
            "budget tight (don't load frontend rules when the agent is reading backend code). Local "
            "for personal preferences without polluting team conventions."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Plan mode vs. headless mode",
        [
            "Plan mode   - interactive, proposes plan, waits for approval (default)",
            "Headless    - claude -p 'prompt' - non-interactive, pipe-friendly",
            "            - --output-format json|stream-json for parsing",
            "            - --allowed-tools 'Read,Edit,Bash' to scope",
            "",
            "Headless is THE CI/CD answer - pre-commit hooks, PR review, release notes",
        ],
        (
            "Three-second demo callout: git log --oneline -20 | claude -p 'summarize these recent "
            "commits'. One command, real value. Headless mode unlocks every CI/CD use case learners "
            "have but didn't know Claude Code supported."
        ),
    ),
    (
        LAYOUT_DEMO,
        "Demo B - CLAUDE.md hierarchy + claude -p (5 min)",
        [
            "Open this repo's own CLAUDE.md - walk the structure",
            "Run: claude -p 'audit this CLAUDE.md against the documented conventions'",
            "Show the headless output",
            "Pipe a git log into it - get a release-notes draft in seconds",
        ],
        (
            "Live demo, 5 minutes. Short because it's the 'oh that's just a CLI' beat - the value is "
            "obvious once they see it. Have a fallback `claude -p` output saved as text in case the "
            "CLI hiccups."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Exercise (5 min)",
        [
            "Design a CLAUDE.md hierarchy for a 3-team monorepo",
            "Teams: backend (Go), frontend (React/TS), infra (Terraform)",
            "Each team has different conventions",
            "Where does each rule go? User, project, subtree, or .claude/rules?",
            "What goes in .claude/rules/ with paths: globs?",
        ],
        (
            "Target answer: project CLAUDE.md for shared (commit style, security), subtree "
            "CLAUDE.md for team-specific conventions, .claude/rules with paths: globs for "
            "language-specific lint rules. Hierarchy is the load-bearing concept."
        ),
    ),
    (
        LAYOUT_SUMMARY,
        "Segment 2 takeaways",
        [
            "Descriptions, not names, drive tool selection",
            ".mcp.json - three transports, ${ENV} expansion, secrets in env never in file",
            "CLAUDE.md hierarchy is additive; claude -p is your CI/CD primitive",
        ],
        (
            "Bridge: 'Tools give the agent hands. Prompts and schemas decide what comes out. "
            "Before we get to structured output in Segment 3, take a 10-minute break - and if "
            "you want a deep dive between segments, Segment 2.5 ships as self-study.'"
        ),
    ),

    # ======================================================================
    # SEGMENT 2.5: BRIDGE (self-study deep dive)
    # ======================================================================
    (
        LAYOUT_UPNEXT,
        "Segment 2.5 - Self-study deep dive (off-clock)",
        [
            "notebooks/segment-2-5-control-surfaces.ipynb",
            "",
            "Three tiers tools come from - server tools, MCP, Claude Code harness",
            "All four tool_choice modes + disable_parallel_tool_use",
            "stop_sequences + max_tokens as control levers + pause_turn + refusal",
            "MCP list_tools runtime discovery",
            "The live Claude Console asset surface - memory_stores, vaults, agents, sessions",
            "",
            "Highest exam-objective density in the repo. Walk it for CCA-F prep.",
        ],
        (
            "Tell learners this is the notebook to re-walk if they're aiming at CCA-F. ~75 minutes "
            "if walked end to end against the live API. Not on the 4-hour clock. The Console asset "
            "surface section is brand new in 2026-05 and won't be in older study guides - this is "
            "where Tim's content has an edge."
        ),
    ),

    # ======================================================================
    # SEGMENT 3: STRUCTURED OUTPUT + RELIABILITY (Domains 4 + 5)
    # ======================================================================
    (
        LAYOUT_SECTION,
        "Segment 3: Structured Output and Reliability",
        None,
        (
            "Domains 4 + 5. 50 minutes. Combined 35% of exam weight. "
            "Live build: invoice extractor with Pydantic + forced tool_choice + validation retry."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Two domains, one segment (again)",
        [
            "Domain 4 - Prompt Engineering and Structured Output (20%)",
            "Domain 5 - Context Management and Reliability (15%)",
            "Combined = 35% of the cert",
            "",
            "Why combined? Structured extraction without context discipline is fragile",
            "Long sessions need context pinning; structured output needs validation retry",
        ],
        (
            "The merge is honest to production. A well-typed extraction call inside a 50-turn "
            "conversation needs both Domain 4 (the call itself) and Domain 5 (context pinning so "
            "turn 47 still has the case facts)."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "What you'll build",
        [
            "An invoice extractor with a Pydantic schema",
            "tool_choice {type: tool, name: extract_invoice} forces valid JSON",
            "Validation-retry loop handles malformed output (max_retries=1)",
            "Three invoices: clean, missing field, ambiguous line item",
            "",
            "By the end of the segment: you can extract typed data from anything",
        ],
        (
            "Set expectations: this is the LONGEST live build (~20 min). Highest demo risk. If "
            "anything goes sideways, switch to the markdown walkthrough next to the notebook. "
            "Sonnet 4.6 for this segment - the one place reasoning depth earns its tokens."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Precise prompts beat vague prompts",
        [
            "Be specific: format, edge cases, missing-data handling",
            "Plausible hallucinations: model invents values for nullable fields",
            "Fix: explicit 'return null if not directly stated' instruction",
            "Inconsistent formats: 'cotton blend' vs 'Cotton/Polyester mix'",
            "Fix: 2-3 few-shot examples, NOT lower temperature",
        ],
        (
            "Few-shot beats temperature tuning every time. The model needs examples, not knobs. "
            "Temperature is the lever you reach for last, not first."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "JSON schema via tool use - the cheat code",
        [
            "1. Define a Pydantic model with your fields",
            "2. Convert to JSON Schema (model.model_json_schema())",
            "3. Register as a tool's input_schema",
            "4. Set tool_choice {type: tool, name: <your_schema>}",
            "5. block.input is now a typed dict matching the schema",
        ],
        (
            "The Pydantic model is the source of truth for both runtime validation AND the API "
            "contract. No JSON parsing, no fence-stripping, no validation rituals. The SDK "
            "validated for you. This is the Domain 4 load-bearing pattern."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Validation + retry loops",
        [
            "Effective: format errors, JSON parse failures, schema mismatches",
            "    -> append the error as a tool_result with is_error: true, retry",
            "    -> succeeds in 2-3 attempts",
            "",
            "INEFFECTIVE: knowledge gaps, ambiguous source data",
            "    -> retries DO NOT help",
            "    -> cap at 2-3, then escalate to human review",
        ],
        (
            "Anti-pattern: infinite retry on missing-info failure. Burns tokens, the model can't "
            "invent the data, escalation is the right move. max_retries=1 is the production default - "
            "one chance to react to its own error, then fail loud."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Confidence-based routing",
        [
            "Have the model output field-level confidence: high | medium | low",
            "Route: high -> automated, medium -> spot-check, low -> human",
            "",
            "CRITICAL: validate accuracy by document type AND field",
            "Aggregate 95% accuracy can mask 60% accuracy on specific segments",
        ],
        (
            "Per-segment calibration warning. If you don't break it down by doc type, "
            "high-aggregate accuracy lies to you. The confidence field is a routing key, not "
            "a vanity metric."
        ),
    ),
    (
        LAYOUT_DEMO,
        "Demo - invoice extractor live build (20 min)",
        [
            "Open notebooks/segment-3-invoice-extractor.ipynb",
            "Build the Pydantic schema, register as a tool",
            "Force tool_choice = {type: tool, name: extract_invoice}",
            "Run on 3 invoices: clean, missing PO, ambiguous total",
            "Add validation-retry loop, watch it recover from a parse failure",
        ],
        (
            "Longest live build of the course. Highest demo risk. Have the markdown walkthrough "
            "open as fallback. Cookbook reference for further study: "
            "claude-cookbooks-main/tool_use/extracting_structured_json.ipynb and "
            "tool_use_with_pydantic.ipynb."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Context preservation in long sessions",
        [
            "Pin case-facts at the top of every turn",
            "Summarize resolved turns into narrative, drop verbatim history",
            "Prune verbose tool outputs - extract the fields you used, drop the rest",
            "Application-side filtering, before the next turn",
            "",
            "How lawyers manage case files. Same discipline for long agent sessions.",
        ],
        (
            "A 48-turn conversation should not have 48 turns of tool_result clutter. The cache "
            "demo in segment-3 shows cache_control on the system block - 2019 tokens cached, "
            "the vendor policy never re-billed across calls."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Escalation design",
        [
            "Honor explicit human requests IMMEDIATELY",
            "    'I want a human NOW' = escalate, don't clarify, don't try one more tool",
            "Pass STRUCTURED summary to the human queue, not the transcript",
            "Never escalate on sentiment alone (frustration != complexity)",
            "Confidence-based escalation: low confidence -> human review",
        ],
        (
            "Frustrated user with a $5 refund doesn't need a human. Calm user with a multi-account "
            "billing dispute does. Sentiment is not complexity. The escalation contract: structured "
            "summary, not the transcript - the human picks up faster, the agent's notes are signal."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Exercise (5 min)",
        [
            "Design an extraction schema for a domain you know",
            "Fields: required vs Optional[]",
            "Add a confidence Literal: 'high' | 'medium' | 'low'",
            "Add a notes Optional[str] for ambiguity",
            "Share your schema in chat",
        ],
        (
            "Reinforces: nullable fields need explicit prompt instructions, confidence is the "
            "routing key, notes catches the 'I'm not sure' cases without crashing validation."
        ),
    ),
    (
        LAYOUT_SUMMARY,
        "Segment 3 takeaways",
        [
            "Pydantic + tool_choice tool = guaranteed valid JSON",
            "Few-shot examples beat temperature tuning every time",
            "Retries help format errors; retries do NOT help missing info",
            "Cache the system block with cache_control - amortize policy across calls",
        ],
        (
            "Bridge: 'You have agents, tools, prompts, validation, context discipline. "
            "Last segment is the cert capstone - what's on the CCA-F exam, how the domain "
            "weights map to today's notebooks, and ten weighted practice questions.'"
        ),
    ),

    # ======================================================================
    # SEGMENT 4: CCA-F CERTIFICATION CAPSTONE
    # ======================================================================
    (
        LAYOUT_SECTION,
        "Segment 4: CCA-F Certification Capstone",
        None,
        (
            "All five CCA-F domains. 50 minutes. 12-minute cert briefing + 28 minutes of "
            "10 weighted practice questions + 10 minutes of Q&A and close."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "What CCA-F is",
        [
            "Claude Certified Architect: Foundations - Anthropic's first official certification",
            "Targets architects shipping with Agent SDK, Claude Code, Claude API, MCP",
            "Roughly the 301 skill level - past 'hello world', not yet 'I run distillation pipelines'",
            "",
            "Beta product as of 2026 - watch the official cert page for changes",
        ],
        (
            "Skills-first framing. You already have the skills. The exam is how you signal them. "
            "This is the moment to mention CCA-F is currently partner-gated; check claude.com/partners "
            "for the current access policy."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Exam mechanics (the load-bearing facts)",
        [
            "Questions: 60 multiple-choice",
            "Time: 120 minutes, single session, no breaks",
            "Scoring: scaled 100-1000, passing = 720, no guess penalty",
            "Delivery: ProctorFree, English",
            "Results: 2 business days",
            "Cost: $99 per attempt",
            "Attempts: ONE - this is the policy that catches people",
        ],
        (
            "The one-attempt rule is the policy nobody plans around. Tell them: don't schedule "
            "until your Anthropic Practice Exam scores are consistently above 900. Treat the $99 "
            "as if it were $990 - because the cost of failing is rescheduling and starting prep over."
        ),
    ),
    (
        LAYOUT_COMPARISON,
        "Domain weights and where each was taught",
        [
            "Domain 1 - Agentic Architecture (27%)",
            "Heaviest. Segment 1 - agents, hooks, coordinator-subagent. Re-walk segment-1.",
            "Domain 2 - Tool Design and MCP (18%)",
            "Segment 2 + 2.5. Tool descriptions, structured errors, MCP list_tools.",
            "Domain 3 - Claude Code (20%)",
            "Segment 2. CLAUDE.md hierarchy, skills, slash commands, claude -p headless.",
            "Domain 4 - Prompts and Structured Output (20%)",
            "Segment 3. Forced tool_choice, Pydantic schemas, few-shot, retry loops.",
            "Domain 5 - Context and Reliability (15%)",
            "Segment 3 + parts of Segment 1. Context pinning, escalation, structured errors.",
        ],
        (
            "D1 + D3 + D4 = 67% of the exam. Weight your study time the same way. D2 + D3 = 38% - "
            "that pairing is why we covered both in one segment. D5 is small but easy to under-teach "
            "because it feels like security common sense until scenario questions get specific."
        ),
    ),
    (
        LAYOUT_DEFINITION,
        "Why the practice sampler is domain-weighted",
        [
            "The naive sampler",
            "Pulls questions evenly across scenario families - which over-tests Claude Code "
            "(two scenarios are heavily D3) and under-tests context management (only fits in one).",
            "Our sampler",
            "Draws 3 D1 + 2 D2 + 2 D3 + 2 D4 + 1 D5 - matches the exam's 27/18/20/20/15 weighting "
            "exactly. Every cohort sees the SAME 10 questions in the same order (seed=2026) for "
            "consistency across deliveries.",
        ],
        (
            "This is why Segment 4 runs against a freshly-tagged practice-questions.json. The "
            "60-question bank distribution doesn't match exam weights, but the SAMPLER does. "
            "Same idea Anthropic uses for the real exam."
        ),
    ),
    (
        LAYOUT_DEMO,
        "Live practice - 10 weighted questions (~28 min)",
        [
            "Open notebooks/segment-4-cca-f-capstone.ipynb",
            "Sampler draws 10 questions by CCA-F domain weight (3/2/2/2/1)",
            "Per question (~2.5 min):",
            "  1. Read the situation",
            "  2. Cohort votes A/B/C/D in chat",
            "  3. Reveal answer + explanation",
            "  4. Color with a production tip from the matching domain-N-*.md",
        ],
        (
            "The community-sourced disclaimer goes here: questions are calibration practice, "
            "NOT exam predictors. Use Anthropic's official Practice Exam as the authoritative "
            "signal before scheduling. Read the disclaimer out loud once at the start."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Tim's week-before-the-exam punchlist",
        [
            "1. Score >900 on the Anthropic Practice Exam cold",
            "2. Rebuild Segment 1 from memory in 30 minutes",
            "3. Re-read domain-1-agentic.md - the 27% domain",
            "4. Skim domain-3-claude-code.md and domain-4-prompts.md - the 20% domains",
            "5. Verify proctoring setup the night before (camera, mic, ID)",
            "6. Single attempt. Do not schedule until 1-4 are green.",
        ],
        (
            "Full 13-item version lives in CERT-PROGRAM-BRIEFING.md. Six headlines on the slide. "
            "Rebuilding Segment 1 from memory in 30 minutes is the load-bearing readiness test - "
            "if you can't, you're not ready to sit the exam."
        ),
    ),

    # ======================================================================
    # COURSE WRAP
    # ======================================================================
    (
        LAYOUT_CONTENT,
        "What you built today",
        [
            "Customer support agent with PreToolUse hook enforcement",
            "Configured Claude Code env with .mcp.json + CLAUDE.md hierarchy",
            "Invoice extractor with Pydantic schema + validation retry",
            "Production reliability mental model - context, escalation, errors, provenance",
            "",
            "Plus self-study: control surfaces and the Console asset surface (Segment 2.5)",
        ],
        (
            "Recap the four live artifacts plus the self-study fifth. Tie back to the cold-open "
            "quote: model is non-deterministic, your architecture cannot be. Every artifact today "
            "is an instance of that pattern."
        ),
    ),
    (
        LAYOUT_REFS,
        "Resources and next steps",
        [
            "Repo",
            "github.com/timothywarner-org/claude-architect",
            "Anthropic cookbook (official, MIT, vendored in this repo)",
            "github.com/anthropics/claude-cookbooks   |   claude-cookbooks-main/",
            "CCA-F study path",
            "EXAM-STUDY-PATH.md - notebook -> CCA-F domain -> scenario mapping",
            "Anthropic Academy",
            "anthropic.skilljar.com - free courses, official Practice Exam",
        ],
        (
            "If they don't ship something in 7 days, the muscle memory fades. Pick the smallest "
            "possible win and ship it. The repo carries EXAM-STUDY-PATH.md which is the post-class "
            "anchor for CCA-F prep - notebook to domain to scenario family in one table."
        ),
    ),
    (
        LAYOUT_CONTENT,
        "Final Q&A (10 min)",
        [
            "How do I monitor an agent in production?",
            "When do I split a monolithic agent into coordinator + subagents?",
            "What is the cost story for prompt caching?",
            "How do I know when I am ready to sit the exam?",
            "What is the renewal path for CCA-F?",
        ],
        (
            "Anticipated questions to invite. Pre-answer the renewal one: Anthropic reserves the "
            "right to retire and refresh exams; certifications tied to beta products may expire "
            "when the production exam goes live. Watch the cert page."
        ),
    ),
    (
        LAYOUT_COVER,
        "Thank you",
        [
            "Tim Warner  |  techtrainertim.com",
            "github.com/timothywarner-org/claude-architect",
            "",
            "Memento mori. Also, ship the PR.",
        ],
        (
            "Final close. The 'memento mori, ship the PR' line is the personal sign-off. Open "
            "for one last round of questions, then wrap. Recording stops after the chat clears."
        ),
    ),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Placeholder indices that the Master 2 layouts use for their PRIMARY body
# content. Different from Master 0 / Master 1 layouts, which use idx 1 or 2.
# We try this list in order when looking for the "main" body placeholder.
BODY_IDX_CANDIDATES = (1, 2, 10, 11, 12, 13, 14, 15, 17, 18, 19)


def find_layout(prs: Presentation, name: str):
    """Return the slide layout matching `name`, preferring the rich library.

    Master 2 in the reference template is the visually-rich Pluralsight-style
    library (Section Header, Demo, Course Title, Important Statement, etc.).
    Master 0 is the basic Office library with the same names but no decorative
    shapes. We iterate from the HIGHEST master index downward so a name that
    appears in both (e.g. "Section Header" exists in M0/L2 and M2/L1) resolves
    to the rich one.
    """
    masters = list(prs.slide_masters)
    for master in reversed(masters):
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(f"Layout not found: {name!r}")


def set_title(slide, text: str) -> bool:
    """Set the slide title via the layout's title placeholder.

    Returns True if a title placeholder was found and populated, False if
    the layout has no title slot. When False, the caller should prepend the
    title to the body content as the first paragraph (in bold) so the title
    still surfaces against the layout's decorative master-level styling.

    No text-box fallback - if the layout doesn't have a title slot, that's
    the layout's decision (Demo, Up Next, Course Title etc. carry their
    branding in master-level decorations and want all text in a single body).
    """
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        return True
    # Some Master 2 layouts use idx=10 as the primary "hero" slot (Course Title).
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10 and ph.has_text_frame:
            ph.text_frame.text = text
            return True
    return False


def set_body(slide, lines: list[str]) -> None:
    """Set the body content via the layout's first body-shaped placeholder.

    Different layouts use different placeholder indices for their primary
    body content. We probe BODY_IDX_CANDIDATES in order, skipping the title
    (idx 0) and any placeholders we've already written into. The first
    non-empty body placeholder gets all the lines as separate paragraphs.

    If the slide has multiple body placeholders (e.g. Comparison with 4),
    we distribute the lines evenly across them. The caller controls the
    distribution by structuring `lines` as alternating label/text pairs.
    """
    # Collect all candidate body placeholders in idx order
    body_phs = []
    for idx_target in BODY_IDX_CANDIDATES:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx_target and ph.has_text_frame:
                if ph not in body_phs:
                    body_phs.append(ph)

    if not body_phs:
        # No usable body placeholder - render as a styled text box fallback.
        # Hopefully unreachable for layouts in this deck; instrumented to
        # surface visual drift if a new layout is added without testing.
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            for run in p.runs:
                run.font.size = Pt(18)
        return

    # Detect the Comparison layout pattern: even number of lines (4, 6, 8) and
    # multiple body placeholders. Distribute as alternating header/text pairs.
    if len(body_phs) >= 4 and len(lines) >= 4 and len(lines) % 2 == 0:
        pairs = list(zip(lines[::2], lines[1::2]))
        for ph, (header, detail) in zip(body_phs, pairs):
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.paragraphs[0].text = header
            p = tf.add_paragraph()
            p.text = detail
        return

    # Default: dump all lines into the first body placeholder.
    target = body_phs[0]
    tf = target.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line


def set_notes(slide, text: str) -> None:
    """Write speaker notes to the slide's notes_slide.

    Paragraph-safe: source text containing `\\n\\n` becomes separate paragraphs
    (visible gap in Camtasia / PowerPoint); single `\\n` becomes a line break.
    Matches the pluralsight-speaker-notes skill convention so future notes
    rewrites stay layout-stable.
    """
    notes_tf = slide.notes_slide.notes_text_frame
    # Clear existing notes content
    paragraphs = list(notes_tf.paragraphs)
    for para in paragraphs[1:]:
        para._p.getparent().remove(para._p)
    first = notes_tf.paragraphs[0]
    for child in list(first._p):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag in ("br", "r"):
            first._p.remove(child)

    # Split on double-newline; each chunk becomes its own paragraph.
    chunks = text.split("\n\n")
    for i, chunk in enumerate(chunks):
        p = first if i == 0 else notes_tf.add_paragraph()
        # Within a chunk, single newlines become line breaks (no paragraph gap).
        lines = chunk.split("\n")
        p.text = lines[0]
        for ln in lines[1:]:
            p.add_run()
            # python-pptx doesn't expose a clean add_line_break; insert via XML.
            from pptx.oxml.ns import qn
            br = p._p.makeelement(qn("a:br"), {})
            p._p.append(br)
            run = p.add_run()
            run.text = ln


def remove_all_slides(prs: Presentation) -> None:
    """Strip every slide from the deck AND drop their XML parts.

    python-pptx's ``slides._sldIdLst.remove()`` leaves orphaned slide{N}.xml
    parts in the package, which trigger duplicate-name warnings on save.
    To truly purge, we drop both the relationship from the presentation part
    and the slide part itself.
    """
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    package = pres_part.package
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        slide_part = pres_part.related_part(rId)
        pres_part.drop_rel(rId)
        sldIdLst.remove(sldId)
        try:
            del package._parts[slide_part.partname]
        except (AttributeError, KeyError):
            pass


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def main() -> int:
    if not REF.exists():
        print(f"ERROR: reference deck missing at {REF}", file=sys.stderr)
        return 1

    prs = Presentation(REF)
    remove_all_slides(prs)

    missing_layouts: set[str] = set()
    fallback_count = 0

    for layout_name, title, body, notes in SLIDES:
        try:
            layout = find_layout(prs, layout_name)
        except KeyError:
            missing_layouts.add(layout_name)
            layout = find_layout(prs, LAYOUT_CONTENT)

        slide = prs.slides.add_slide(layout)
        had_title_slot = set_title(slide, title)
        # If the layout has no title placeholder, fold the title into the body
        # so the master-level decoration carries the visual and the title text
        # still appears (as the first paragraph). The Demo / Up Next / Course
        # Title / Module Overview layouts in Master 2 work this way.
        effective_body = body
        if not had_title_slot and title:
            effective_body = [title] + ([""] if body else []) + (body or [])
        if effective_body:
            set_body(slide, effective_body)
        if notes:
            set_notes(slide, notes)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)

    size_mb = OUT.stat().st_size / (1024 * 1024)
    print(f"Wrote {OUT}")
    print(f"  Slides: {len(SLIDES)}")
    print(f"  Size:   {size_mb:.2f} MB")
    if missing_layouts:
        print(f"  WARN: layouts not found, fell back to {LAYOUT_CONTENT}: {sorted(missing_layouts)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
