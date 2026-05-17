"""Inspect a reference PPTX deck to extract layout names, theme, and slide structure.

Outputs a JSON-like dict the deck builder can consume to inherit the look-and-feel
of an existing O'Reilly deck (in this case, context-engineering-april-2026.pptx)
when generating a new course deck.

CCA-F course author convention: every script self-documents its purpose.
"""

import json
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

REF = Path(r"C:/github/context-engineering/instructor/context-engineering-april-2026.pptx")


def emu_to_in(emu: int) -> float:
    """EMUs to inches (914400 per inch)."""
    return round(emu / 914400, 2)


def inspect(prs: Presentation) -> dict:
    out = {
        "slide_count": len(prs.slides),
        "slide_width_in": emu_to_in(prs.slide_width),
        "slide_height_in": emu_to_in(prs.slide_height),
        "masters": [],
        "layouts": [],
        "layout_usage": Counter(),
        "first_10_slides": [],
        "section_dividers": [],
        "speaker_notes_present_count": 0,
    }

    for m_idx, master in enumerate(prs.slide_masters):
        out["masters"].append(
            {
                "index": m_idx,
                "name": master.name,
                "layout_count": len(master.slide_layouts),
            }
        )
        for l_idx, layout in enumerate(master.slide_layouts):
            out["layouts"].append(
                {
                    "master_index": m_idx,
                    "layout_index": l_idx,
                    "name": layout.name,
                    "placeholder_names": [ph.name for ph in layout.placeholders],
                }
            )

    layout_names_by_id = {id(l): l.name for m in prs.slide_masters for l in m.slide_layouts}

    for s_idx, slide in enumerate(prs.slides):
        layout_name = layout_names_by_id.get(id(slide.slide_layout), "UNKNOWN")
        out["layout_usage"][layout_name] += 1

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            out["speaker_notes_present_count"] += 1

        if s_idx < 10:
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().split("\n")[0][:80]
                    break
            out["first_10_slides"].append(
                {
                    "index": s_idx,
                    "layout": layout_name,
                    "title_preview": title,
                }
            )

        # Section dividers: layouts whose name contains "section" or "divider"
        if "section" in layout_name.lower() or "divider" in layout_name.lower():
            out["section_dividers"].append({"index": s_idx, "layout": layout_name})

    out["layout_usage"] = dict(out["layout_usage"].most_common())
    return out


def main() -> int:
    if not REF.exists():
        print(f"ERROR: reference deck not found at {REF}", file=sys.stderr)
        return 1

    prs = Presentation(REF)
    report = inspect(prs)
    print(json.dumps(report, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
